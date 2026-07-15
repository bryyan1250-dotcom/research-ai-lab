from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path.cwd()
OUTPUT_DIR = ROOT / "outputs"
PACKAGE_DIR = ROOT / "launch-keynote"

ACCENT = RGBColor(180, 35, 24)
TEXT = RGBColor(16, 24, 40)
MUTED = RGBColor(102, 112, 133)
FAINT = RGBColor(208, 213, 221)
SOFT = RGBColor(254, 228, 226)


tool_slides = [
    {
        "tool": "Perplexity",
        "mission": "Entender el problema",
        "open": "perplexity.ai",
        "search": "How is Artificial Intelligence impacting productivity and economic growth in emerging economies?",
        "observe": "Conceptos, reportes, instituciones",
        "output": "Mapa inicial de palabras clave",
        "next": "OpenAlex",
    },
    {
        "tool": "OpenAlex",
        "mission": "Abrir el universo académico",
        "open": "openalex.org",
        "search": "\"artificial intelligence\" productivity \"emerging economies\"",
        "observe": "Autores, revistas, años, volumen",
        "output": "Primer corpus académico",
        "next": "Semantic Scholar",
    },
    {
        "tool": "Semantic Scholar",
        "mission": "Refinar relevancia",
        "open": "semanticscholar.org",
        "search": "\"artificial intelligence\" productivity \"emerging economies\"",
        "observe": "Trabajos clave y citaciones",
        "output": "Lista de lectura priorizada",
        "next": "ResearchRabbit",
    },
    {
        "tool": "ResearchRabbit",
        "mission": "Seguir conexiones",
        "open": "researchrabbit.ai",
        "search": "AI productivity emerging economies",
        "observe": "Redes, clústeres, autores",
        "output": "Grafo de literatura",
        "next": "Litmaps",
    },
    {
        "tool": "Litmaps",
        "mission": "Ver evolución",
        "open": "litmaps.com",
        "search": "AI productivity emerging economies",
        "observe": "Huecos temporales y ramas",
        "output": "Mapa cronológico útil",
        "next": "Consensus",
    },
    {
        "tool": "Consensus",
        "mission": "Escuchar la evidencia",
        "open": "consensus.app",
        "search": "Does artificial intelligence increase productivity in developing countries?",
        "observe": "Patrones y convergencias",
        "output": "Dirección general del debate",
        "next": "Scite",
    },
    {
        "tool": "Scite",
        "mission": "Validar confianza",
        "open": "scite.ai",
        "search": "\"artificial intelligence\" productivity firms growth \"developing countries\"",
        "observe": "Supporting, contrasting, mentioning",
        "output": "Evidencia más robusta",
        "next": "SciSpace",
    },
    {
        "tool": "SciSpace",
        "mission": "Leer con precisión",
        "open": "typeset.io",
        "search": "\"artificial intelligence\" productivity \"developing economies\"",
        "observe": "Métodos, datos, limitaciones",
        "output": "Notas analíticas comparables",
        "next": "NotebookLM",
    },
    {
        "tool": "NotebookLM",
        "mission": "Construir memoria",
        "open": "notebooklm.google.com",
        "search": "Corpus: 60 artículos validados",
        "observe": "Temas, tensiones, síntesis",
        "output": "Base de conocimiento compartida",
        "next": "ChatGPT",
    },
    {
        "tool": "ChatGPT",
        "mission": "Comunicar ciencia",
        "open": "chatgpt.com",
        "search": "Synthesize validated evidence. Compare findings, controversies, and one research gap.",
        "observe": "Argumento, estructura, claridad",
        "output": "Síntesis científica lista",
        "next": "Nueva oportunidad",
    },
]


slides = [
    {
        "number": 1,
        "key": "cover",
        "chapter": "INTRO",
        "title": "Research AI Lab",
        "subtitle": "Versión 1.0",
        "kicker": "Open Research. Open Science. Open Intelligence.",
        "note": [
            "Abrir con calma.",
            "Presentar el lanzamiento como el inicio de una iniciativa, no como una demo.",
            "Marcar que hoy vamos a recorrer una investigación completa.",
        ],
    },
    {
        "number": 2,
        "key": "story",
        "chapter": "WHY",
        "title": "Imagina que mañana debes empezar tu estado del arte.",
        "subtitle": "Sin autores. Sin papers. Sin mapa.",
        "note": [
            "Bajar la voz.",
            "Invitar a la audiencia a entrar en la escena de una investigadora real.",
            "Generar tensión antes de mencionar cualquier herramienta.",
        ],
    },
    {
        "number": 3,
        "key": "quote",
        "chapter": "WHY",
        "title": '"Prepara un estado del arte completo sobre este tema."',
        "subtitle": "Solo tienes una pregunta.",
        "note": [
            "Hacer una pausa antes y después de la cita.",
            "Mirar a la audiencia.",
            "Dejar que el problema se sienta grande.",
        ],
    },
    {
        "number": 4,
        "key": "question",
        "chapter": "WHY",
        "title": "¿Cómo está impactando la inteligencia artificial la productividad y el crecimiento económico en economías emergentes?",
        "subtitle": "",
        "note": [
            "Esta es la pregunta central del laboratorio.",
            "No explicar todavía cómo se resolverá.",
            "Solo dejarla suspendida.",
        ],
    },
    {
        "number": 5,
        "key": "statement",
        "chapter": "WHY",
        "title": "Toda investigación comienza con incertidumbre.",
        "subtitle": "La diferencia está en cómo la orquestamos.",
        "note": [
            "Conectar el problema individual con el problema estructural de la investigación moderna.",
            "Preparar la transición hacia el capítulo WHAT.",
        ],
    },
    {
        "number": 6,
        "key": "chapter",
        "chapter": "WHAT",
        "title": "WHAT",
        "subtitle": "Por qué el flujo tradicional ya no escala.",
        "note": [
            "Cambiar energía.",
            "Pasar de la historia personal al diagnóstico del sistema.",
        ],
    },
    {
        "number": 7,
        "key": "verticalFlow",
        "chapter": "WHAT",
        "title": "Flujo tradicional",
        "subtitle": "Buscar. Descargar. Copiar. Perder contexto.",
        "flow": ["Google", "Google Scholar", "PDFs", "Excel", "Word", "Publicación"],
        "note": [
            "Recorrer el flujo lentamente.",
            "Mostrar que no es un problema de esfuerzo, sino de escala y fragmentación.",
        ],
    },
    {
        "number": 8,
        "key": "statement",
        "chapter": "WHAT",
        "title": "Investigar ya no es solo buscar.",
        "subtitle": "Investigar es orquestar conocimiento.",
        "note": [
            "Esta es la tesis conceptual del keynote.",
            "Decirla como una definición, no como un eslogan.",
        ],
    },
    {
        "number": 9,
        "key": "verticalFlow",
        "chapter": "WHAT",
        "title": "Research AI Stack",
        "subtitle": "Descubrimiento. lectura. evidencia. memoria. escritura.",
        "flow": ["Descubrimiento", "Expansión", "Validación", "Lectura", "Memoria", "Comunicación"],
        "note": [
            "Presentar el stack como una arquitectura mental.",
            "Todavía no entrar en productos concretos.",
        ],
    },
    {
        "number": 10,
        "key": "chapter",
        "chapter": "HOW",
        "title": "HOW",
        "subtitle": "Laboratorio 01: Economía.",
        "note": [
            "Marcar que ahora la audiencia deja de ser espectadora y se convierte en investigadora.",
        ],
    },
    {
        "number": 11,
        "key": "story",
        "chapter": "HOW",
        "title": "Research AI Lab",
        "subtitle": "Laboratorio 01: Economía",
        "note": [
            "Presentar el laboratorio como el ejemplo canónico de la metodología.",
            "No venderlo como curso ni como colección de prompts.",
        ],
    },
    {
        "number": 12,
        "key": "workflow",
        "chapter": "HOW",
        "title": "Del problema a la oportunidad científica.",
        "subtitle": "",
        "workflow": [
            "Pregunta",
            "Perplexity",
            "OpenAlex",
            "Semantic Scholar",
            "ResearchRabbit",
            "Litmaps",
            "Consensus",
            "Scite",
            "SciSpace",
            "NotebookLM",
            "ChatGPT",
            "Nueva oportunidad",
        ],
        "note": [
            "Este es el esqueleto del keynote.",
            "Explicar que cada herramienta cumple una misión distinta dentro de una sola investigación.",
        ],
    },
]

for i, item in enumerate(tool_slides, start=13):
    slides.append(
        {
            "number": i,
            "key": "tool",
            "chapter": "HOW",
            "title": item["tool"],
            "note": [
                f"Presentar {item['tool']} como una decisión metodológica.",
                "No explicar funciones genéricas.",
                "Cerrar mostrando cómo alimenta el siguiente paso.",
            ],
            **item,
        }
    )

slides.extend(
    [
        {
            "number": 23,
            "key": "chapter",
            "chapter": "RESULT",
            "title": "RESULT",
            "subtitle": "La transformación de una pregunta en conocimiento estructurado.",
            "note": [
                "Subir el ritmo.",
                "Mostrar que el laboratorio realmente produjo conocimiento nuevo.",
            ],
        },
        {
            "number": 24,
            "key": "verticalFlow",
            "chapter": "RESULT",
            "title": "Transformación",
            "subtitle": "",
            "flow": [
                "1200 papers",
                "240 relevantes",
                "60 sólidos",
                "12 revisiones",
                "5 autores clave",
                "3 escuelas",
                "2 controversias",
                "1 vacío",
            ],
            "note": [
                "Recorrer el embudo como una prueba de que la metodología reduce complejidad.",
                "Mantenerlo visual y memorable.",
            ],
        },
        {
            "number": 25,
            "key": "statement",
            "chapter": "RESULT",
            "title": "Toda buena investigación termina con una nueva pregunta.",
            "subtitle": "Ahí empieza la siguiente oportunidad científica.",
            "note": [
                "Resolver el arco narrativo.",
                "Mostrar que el valor no es solo resumir, sino abrir un nuevo frente de investigación.",
            ],
        },
        {
            "number": 26,
            "key": "chapter",
            "chapter": "WHAT NEXT",
            "title": "WHAT NEXT",
            "subtitle": "El lanzamiento oficial de Research AI Lab.",
            "note": [
                "Pasar del laboratorio al movimiento.",
                "Decir explícitamente que la presentación no es el producto.",
            ],
        },
        {
            "number": 27,
            "key": "launch",
            "chapter": "WHAT NEXT",
            "title": "Versión 1.0",
            "subtitle": "Arquitectura pública para investigar mejor con inteligencia artificial.",
            "items": [
                "Laboratorio de Economía",
                "Research AI Toolkit",
                "Research Recipes",
                "Resources",
                "GitHub Repository",
            ],
            "note": [
                "Presentar cada módulo como parte de una arquitectura coherente.",
                "Evitar sonar como lista de funcionalidades.",
            ],
        },
        {
            "number": 28,
            "key": "roadmap",
            "chapter": "WHAT NEXT",
            "title": "Roadmap",
            "subtitle": "",
            "roadmap": [
                ["Versión 1.1", "Education", "Coming Soon"],
                ["Versión 1.2", "Artificial Intelligence", "Coming Soon"],
                ["Versión 2", "Research Agents", ""],
                ["Versión 3", "Research Operating System", ""],
            ],
            "note": [
                "Mostrar ambición con estructura.",
                "El mensaje es continuidad, no improvisación.",
            ],
        },
        {
            "number": 29,
            "key": "ecosystem",
            "chapter": "WHAT NEXT",
            "title": "La investigadora sigue en el centro.",
            "subtitle": "Descubrimiento. lectura. evidencia. citación. memoria. escritura. monitoreo. agentes.",
            "note": [
                "Reforzar la idea central: la inteligencia artificial no reemplaza a la investigadora.",
                "La conecta con un ecosistema especializado.",
            ],
        },
        {
            "number": 30,
            "key": "qr",
            "chapter": "WHAT NEXT",
            "title": "Continúa el laboratorio.",
            "subtitle": "Research AI Lab · Versión 1.0",
            "note": [
                "Cerrar con llamada a la acción.",
                "Explicar que el QR debe apuntar al repositorio público oficial una vez se publique.",
            ],
        },
        {
            "number": 31,
            "key": "backup",
            "chapter": "BACKUP",
            "title": "Backup · Stack completo",
            "subtitle": "Qué problema resuelve cada capa metodológica.",
            "backup_lines": [
                "Descubrimiento · comprender y expandir",
                "Validación · medir confianza",
                "Lectura · comparar método y datos",
                "Memoria · sostener contexto",
                "Comunicación · escribir con evidencia",
            ],
            "note": [
                "Solo usar si alguien pide una explicación más explícita del stack.",
            ],
        },
        {
            "number": 32,
            "key": "backup",
            "chapter": "BACKUP",
            "title": "Backup · Consultas exactas",
            "subtitle": "Punto de partida para repetir el laboratorio.",
            "backup_lines": [
                "\"artificial intelligence\" productivity \"emerging economies\"",
                "\"AI adoption\" firms productivity growth",
                "\"developing countries\" artificial intelligence growth",
            ],
            "note": [
                "Solo mostrar si la audiencia pide replicabilidad operativa inmediata.",
            ],
        },
        {
            "number": 33,
            "key": "backup",
            "chapter": "BACKUP",
            "title": "Backup · Checklist en vivo",
            "subtitle": "Antes de abrir herramientas en escenario.",
            "backup_lines": [
                "Navegadores limpios",
                "Sesiones abiertas",
                "Pestañas fijadas",
                "Corpus listo",
                "Plan B local",
            ],
            "note": [
                "Utilizar antes del evento, o si la audiencia pregunta por la operación en vivo.",
            ],
        },
    ]
)


def inches(value: float) -> float:
    return Inches(value)


def add_textbox(slide, x, y, w, h, text, size=24, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_rule(slide, x, y, w, h=0.03, color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, x, y, size, fill=SOFT, line=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, inches(x), inches(y), inches(size), inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_chapter_header(slide, chapter):
    add_textbox(slide, 0.95, 0.22, 2.1, 0.24, chapter, size=11, bold=True, color=MUTED)


def draw_cover(slide, item):
    add_textbox(slide, 0.95, 0.28, 2.0, 0.24, "Launch Keynote", size=11, bold=True, color=MUTED)
    add_rule(slide, 0.95, 1.22, 1.0, 0.04)
    add_textbox(slide, 0.95, 1.52, 7.0, 1.1, item["title"], size=46, bold=True)
    add_textbox(slide, 0.95, 2.75, 3.0, 0.45, item["subtitle"], size=22, bold=True, color=ACCENT)
    add_textbox(slide, 0.95, 3.55, 7.8, 0.55, item["kicker"], size=20, color=MUTED)


def draw_story(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_rule(slide, 0.95, 1.25, 0.8, 0.04)
    add_textbox(slide, 0.95, 1.52, 10.0, 1.35, item["title"], size=40, bold=True)
    if item["subtitle"]:
        add_textbox(slide, 0.95, 3.12, 8.0, 0.55, item["subtitle"], size=22, color=MUTED)


def draw_quote(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 1.7, 1.95, 10.0, 1.75, item["title"], size=36, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 3.6, 4.18, 6.0, 0.45, item["subtitle"], size=20, color=ACCENT, align=PP_ALIGN.CENTER)


def draw_question(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 1.25, 1.55, 10.8, 3.3, item["title"], size=34, bold=True, align=PP_ALIGN.CENTER)


def draw_statement(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 1.15, 1.95, 11.0, 1.3, item["title"], size=42, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.55, 3.7, 8.2, 0.55, item["subtitle"], size=20, color=MUTED, align=PP_ALIGN.CENTER)


def draw_chapter(slide, item):
    add_rule(slide, 0.95, 1.12, 1.0, 0.04)
    add_textbox(slide, 0.95, 1.62, 6.0, 1.2, item["title"], size=62, bold=True)
    add_textbox(slide, 0.95, 3.05, 8.0, 0.55, item["subtitle"], size=20, color=MUTED)


def draw_vertical_flow(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.72, 4.8, 0.45, item["title"], size=28, bold=True)
    if item["subtitle"]:
        add_textbox(slide, 0.95, 1.18, 7.2, 0.35, item["subtitle"], size=15, color=MUTED)
    start_y = 1.9
    for idx, line in enumerate(item["flow"]):
        color = ACCENT if idx == len(item["flow"]) - 1 else TEXT
        add_textbox(slide, 4.25, start_y + idx * 0.62, 4.3, 0.35, line, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        if idx < len(item["flow"]) - 1:
            add_textbox(slide, 6.15, start_y + idx * 0.62 + 0.26, 0.45, 0.25, "↓", size=18, color=MUTED, align=PP_ALIGN.CENTER)


def draw_workflow(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.72, 8.0, 0.45, item["title"], size=28, bold=True)
    row1 = item["workflow"][:6]
    row2 = item["workflow"][6:]
    for idx, label in enumerate(row1):
        x = 0.6 + idx * 2.0
        color = ACCENT if idx == 0 else TEXT
        add_textbox(slide, x, 2.18, 1.5, 0.6, label, size=17, bold=True, color=color, align=PP_ALIGN.CENTER)
        if idx < len(row1) - 1:
            add_textbox(slide, x + 1.45, 2.32, 0.3, 0.2, "→", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    for idx, label in enumerate(row2):
        x = 0.6 + idx * 2.0
        color = ACCENT if label == "Nueva oportunidad" else TEXT
        add_textbox(slide, x, 4.2, 1.5, 0.6, label, size=17, bold=True, color=color, align=PP_ALIGN.CENTER)
        if idx < len(row2) - 1:
            add_textbox(slide, x + 1.45, 4.34, 0.3, 0.2, "→", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 11.7, 3.08, 0.3, 0.2, "↓", size=16, color=MUTED, align=PP_ALIGN.CENTER)


def draw_tool(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.62, 3.4, 0.45, item["tool"], size=32, bold=True)
    add_rule(slide, 0.95, 1.12, 0.8, 0.03)

    def label(x, y, text):
        add_textbox(slide, x, y, 1.2, 0.2, text, size=11, bold=True, color=MUTED)

    label(0.95, 1.52, "MISSION")
    add_textbox(slide, 0.95, 1.78, 3.8, 0.75, item["mission"], size=20, bold=True)
    label(0.95, 2.82, "OPEN")
    add_textbox(slide, 0.95, 3.08, 3.8, 0.35, item["open"], size=18, bold=True, color=ACCENT)
    label(0.95, 4.12, "OUTPUT")
    add_textbox(slide, 0.95, 4.38, 3.8, 0.6, item["output"], size=20, bold=True)
    label(0.95, 5.25, "NEXT")
    add_textbox(slide, 0.95, 5.48, 2.5, 0.35, item["next"], size=18, bold=True, color=ACCENT)

    label(6.55, 1.52, "SEARCH")
    add_textbox(slide, 6.55, 1.78, 5.0, 1.2, item["search"], size=15)
    label(6.55, 4.02, "OBSERVE")
    add_textbox(slide, 6.55, 4.28, 4.5, 0.55, item["observe"], size=17)


def draw_launch(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.7, 4.0, 0.45, item["title"], size=30, bold=True)
    add_textbox(slide, 0.95, 1.18, 9.0, 0.35, item["subtitle"], size=16, color=MUTED)
    for idx, text in enumerate(item["items"]):
        col = idx % 2
        row = idx // 2
        x = 0.95 + col * 4.6
        y = 2.4 + row * 1.25
        add_rule(slide, x, y, 0.7, 0.03)
        add_textbox(slide, x, y + 0.18, 3.7, 0.5, text, size=22, bold=True)


def draw_roadmap(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.7, 3.0, 0.45, item["title"], size=30, bold=True)
    for idx, (version, name, status) in enumerate(item["roadmap"]):
        x = 0.95 + idx * 2.6
        add_rule(slide, x, 2.2, 0.7, 0.03)
        add_textbox(slide, x, 2.45, 2.1, 0.3, version, size=16, bold=True, color=ACCENT)
        add_textbox(slide, x, 2.92, 2.2, 0.95, name, size=22, bold=True)
        if status:
            add_textbox(slide, x, 4.08, 2.1, 0.25, status, size=14, color=MUTED)


def draw_ecosystem(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.62, 5.8, 0.45, item["title"], size=28, bold=True)
    add_circle(slide, 5.2, 2.6, 2.1)
    add_textbox(slide, 5.45, 3.45, 1.6, 0.3, "Investigadora", size=20, bold=True, align=PP_ALIGN.CENTER)
    labels = [
        ("Discovery", 1.45, 2.0),
        ("Reading", 10.15, 2.0),
        ("Evidence", 10.7, 3.35),
        ("Citation", 9.95, 5.05),
        ("Writing", 1.55, 5.05),
        ("Workspace", 1.0, 3.35),
        ("Monitoring", 3.15, 5.55),
        ("Agents", 8.4, 5.55),
        ("Knowledge Graph", 4.25, 1.35),
    ]
    for label, x, y in labels:
        add_circle(slide, x, y, 0.82, RGBColor(255, 255, 255), FAINT)
        add_textbox(slide, x + 0.03, y + 0.25, 0.76, 0.2, label, size=11, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        2.55,
        6.45,
        8.2,
        0.42,
        "La inteligencia artificial no reemplaza a la investigadora. Amplifica su capacidad de crear conocimiento.",
        size=15,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def draw_qr(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.78, 5.0, 0.45, item["title"], size=34, bold=True)
    add_textbox(slide, 0.95, 1.28, 4.0, 0.28, item["subtitle"], size=16, color=ACCENT)
    qr = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inches(8.75), inches(1.9), inches(2.3), inches(2.3)
    )
    qr.fill.solid()
    qr.fill.fore_color.rgb = RGBColor(255, 255, 255)
    qr.line.color.rgb = TEXT
    add_textbox(slide, 9.05, 2.7, 1.75, 0.55, "QR oficial\ndel proyecto", size=19, bold=True, align=PP_ALIGN.CENTER)
    for idx, label in enumerate(["GitHub", "Toolkit", "Economics Lab", "Research Recipes", "Roadmap"]):
        y = 2.75 + idx * 0.65
        add_rule(slide, 0.95, y, 0.55, 0.03)
        add_textbox(slide, 0.95, y + 0.18, 3.3, 0.25, label, size=19, bold=True)


def draw_backup(slide, item):
    add_chapter_header(slide, item["chapter"])
    add_textbox(slide, 0.95, 0.72, 5.2, 0.45, item["title"], size=28, bold=True)
    add_textbox(slide, 0.95, 1.18, 7.5, 0.28, item["subtitle"], size=15, color=MUTED)
    for idx, line in enumerate(item["backup_lines"]):
        y = 2.3 + idx * 0.92
        add_rule(slide, 0.95, y, 0.52, 0.03)
        add_textbox(slide, 0.95, y + 0.18, 9.6, 0.28, line, size=19, bold=True)


def render_slide(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if item["key"] == "cover":
        draw_cover(slide, item)
    elif item["key"] == "story":
        draw_story(slide, item)
    elif item["key"] == "quote":
        draw_quote(slide, item)
    elif item["key"] == "question":
        draw_question(slide, item)
    elif item["key"] == "statement":
        draw_statement(slide, item)
    elif item["key"] == "chapter":
        draw_chapter(slide, item)
    elif item["key"] == "verticalFlow":
        draw_vertical_flow(slide, item)
    elif item["key"] == "workflow":
        draw_workflow(slide, item)
    elif item["key"] == "tool":
        draw_tool(slide, item)
    elif item["key"] == "launch":
        draw_launch(slide, item)
    elif item["key"] == "roadmap":
        draw_roadmap(slide, item)
    elif item["key"] == "ecosystem":
        draw_ecosystem(slide, item)
    elif item["key"] == "qr":
        draw_qr(slide, item)
    elif item["key"] == "backup":
        draw_backup(slide, item)


def story_objective(item):
    if item["key"] == "tool":
        return f"Explicar la misión metodológica de {item['tool']} dentro del flujo de investigación."
    mapping = {
        "cover": "Presentar el lanzamiento como el nacimiento de una iniciativa abierta.",
        "story": "Abrir una escena narrativa o introducir el laboratorio como experiencia.",
        "quote": "Materializar la presión real que enfrenta una investigadora.",
        "question": "Fijar la pregunta canónica del laboratorio.",
        "statement": "Condensar la tesis conceptual de la sección.",
        "chapter": "Separar capítulos y resetear la atención del público.",
        "verticalFlow": "Visualizar un flujo o una transformación secuencial.",
        "workflow": "Mostrar la arquitectura completa del laboratorio.",
        "launch": "Lanzar públicamente los componentes de la Versión 1.0.",
        "roadmap": "Demostrar continuidad estratégica y expansión futura.",
        "ecosystem": "Cerrar con una visión humana y conectada del futuro de la investigación.",
        "qr": "Invitar a la audiencia a continuar el laboratorio fuera del escenario.",
        "backup": "Respaldar preguntas operativas o metodológicas en Q&A.",
    }
    return mapping[item["key"]]


def visual_treatment(item):
    mapping = {
        "cover": "Portada minimalista con regla roja, tipografía dominante y mucho espacio en blanco.",
        "story": "Tipografía grande, centro de gravedad claro y composición sobria.",
        "quote": "Tipografía grande, centro de gravedad claro y composición sobria.",
        "question": "Tipografía grande, centro de gravedad claro y composición sobria.",
        "statement": "Tipografía grande, centro de gravedad claro y composición sobria.",
        "chapter": "Tipografía grande, centro de gravedad claro y composición sobria.",
        "verticalFlow": "Secuencia vertical de lectura inmediata con flechas mínimas.",
        "workflow": "Mapa horizontal en dos filas para sostener la narrativa del laboratorio.",
        "tool": "Plantilla fija de dos columnas con jerarquía editorial consistente.",
        "launch": "Lista estructurada con módulos grandes y respiración amplia.",
        "roadmap": "Cuatro columnas limpias con énfasis en versión y nombre.",
        "ecosystem": "Composición radial alrededor de la investigadora.",
        "qr": "Llamado a la acción con placeholder funcional para QR.",
        "backup": "Slides utilitarias, legibles y visualmente consistentes con la cubierta principal.",
    }
    return mapping[item["key"]]


def script_for(item):
    if item["key"] == "cover":
        return (
            "Hoy no vengo a presentar una herramienta. Vengo a lanzar una iniciativa abierta para investigar "
            "mejor con inteligencia artificial. Research AI Lab nace como una arquitectura pública para organizar "
            "conocimiento, sostener rigor y construir una nueva práctica de investigación."
        )
    if item["key"] == "question":
        return (
            "No empecemos por la herramienta. Empecemos por la pregunta. Esta pregunta será nuestro laboratorio "
            "en vivo y el hilo conductor de todo el keynote."
        )
    if item["key"] == "workflow":
        return (
            "Este es el recorrido completo. No son aplicaciones aisladas. Es una secuencia metodológica. "
            "Cada paso reduce incertidumbre y prepara el siguiente."
        )
    if item["key"] == "launch":
        return (
            "Esto es lo que hoy se lanza públicamente. No solo una presentación, sino una base abierta para "
            "laboratorios, toolkit, recetas y recursos de investigación."
        )
    if item["key"] == "ecosystem":
        return (
            "La inteligencia artificial no reemplaza a la investigadora. La conecta con una red de inteligencias "
            "especializadas. El centro sigue siendo humano. Lo que cambia es la capacidad de orquestación."
        )
    if item["key"] == "qr":
        return (
            "La presentación termina aquí, pero el laboratorio continúa. El QR debe abrir el repositorio oficial "
            "y convertirse en el punto de entrada real a la iniciativa."
        )
    if item["key"] == "tool":
        return (
            f"Nuestra misión en {item['tool']} no es explorar funciones. Es tomar una decisión metodológica concreta: "
            f"{item['mission'].lower()}. Lo importante es qué observamos, qué extraemos y cómo ese resultado alimenta "
            "el siguiente paso."
        )
    return " ".join(item["note"])


def write_markdown(name: str, content: str):
    (PACKAGE_DIR / name).write_text(content, encoding="utf-8")


def build_docs():
    storyboard = ["# Storyboard", "", "## Estructura general", ""]
    script = ["# Presenter Script", ""]
    notes = ["# Speaker Notes", ""]

    for item in slides:
        storyboard.extend(
            [
                f"## Slide {item['number']} · {item['title']}",
                f"Capítulo: {item['chapter']}",
                f"Objetivo: {story_objective(item)}",
                f"Tratamiento visual: {visual_treatment(item)}",
                "",
            ]
        )
        script.extend([f"## Slide {item['number']} · {item['title']}", script_for(item), ""])
        notes.extend([f"## Slide {item['number']} · {item['title']}"])
        notes.extend([f"- {line}" for line in item["note"]])
        notes.append("")

    write_markdown("storyboard.md", "\n".join(storyboard) + "\n")
    write_markdown("presenter-script.md", "\n".join(script) + "\n")
    write_markdown("speaker-notes.md", "\n".join(notes) + "\n")
    write_markdown(
        "animation-suggestions.md",
        """# Animation Suggestions

## Principios

- Mantener animaciones discretas, limpias y consistentes.
- Usar aparición simple o desvanecimiento corto.
- Evitar efectos teatrales, rebotes o zooms agresivos.

## Recomendaciones por tipo de slide

- Portada: aparición del título y luego subtítulo.
- Slides narrativas: fade simple del bloque principal.
- Flujos verticales: revelar un paso a la vez.
- Workflow completo: revelar por tramos, no todos los nodos a la vez.
- Tool slides: mantener fija la plantilla y animar solo el contenido principal.
- Roadmap: entrada por columnas, de izquierda a derecha.
- Ecosistema final: primero la investigadora, luego los nodos periféricos.
- QR final: revelar el bloque QR después de la frase de cierre.
""",
    )
    write_markdown(
        "timing-45-minutes.md",
        """# Timing · 45 Minutes

## Propuesta

- Slides 1-5: 7 minutos
- Slides 6-9: 6 minutos
- Slides 10-12: 5 minutos
- Slides 13-22: 16 minutos
- Slides 23-25: 4 minutos
- Slides 26-30: 5 minutos
- Q&A o transición a demo: 2 minutos

## Ritmo recomendado

- Abrir lento y con pausa.
- Acelerar en WHAT.
- Sostener el centro del keynote en HOW.
- Cerrar con energía institucional, no comercial.
""",
    )
    write_markdown(
        "timing-20-minutes.md",
        """# Timing · 20 Minutes

## Corte recomendado

- Slides 1-5: 4 minutos
- Slides 6-9: 3 minutos
- Slides 10-12: 2 minutos
- Slides 13, 14, 16, 18, 21, 22: 6 minutos
- Slides 24-30: 5 minutos

## Regla de compresión

- Mostrar solo seis tool slides.
- Mantener la lógica completa del flujo.
- Usar backup slides solo en preguntas.
""",
    )
    write_markdown(
        "live-lab-checklist.md",
        """# Live Laboratory Checklist

## Antes del evento

- Verificar conexión estable.
- Abrir sesión en todas las herramientas.
- Preparar pestañas en el orden del flujo.
- Tener corpus local de respaldo.
- Confirmar fuente y tamaño de pantalla del venue.

## Antes de subir al escenario

- Cerrar notificaciones.
- Limpiar historial reciente visible.
- Activar modo no molestar.
- Tener la presentación y el navegador en ventanas separadas.
- Confirmar acceso rápido al slide 12 y al slide 30.

## Plan B

- Si falla una herramienta, pasar al slide siguiente y explicar la lógica metodológica.
- Si falla internet, usar slides de backup y corpus local.
- Si el tiempo se comprime, saltar a las slides 24, 27, 28, 29 y 30.
""",
    )
    write_markdown(
        "backup-slides.md",
        """# Backup Slides

## Slide 31

- Explica el stack completo por capas metodológicas.

## Slide 32

- Ofrece consultas exactas para replicar el laboratorio.

## Slide 33

- Resume el checklist operativo del laboratorio en vivo.
""",
    )
    write_markdown(
        "visual-identity-review.md",
        """# Visual Identity Consistency Review

## Lo que se mantiene consistente

- Fondo blanco en todo el deck.
- Tipografía dominante y de gran escala.
- Uso disciplinado de un único acento rojo.
- Espacio en blanco amplio.
- Ausencia de tarjetas densas y estética de interfaz.

## Riesgos a evitar

- Cargar demasiado texto en las tool slides.
- Reducir tipografía para hacer entrar contenido.
- Convertir el keynote en tutorial de producto.
- Mezclar estilos de iconografía o ilustración de terceros.

## Revisión final recomendada

- Verificar que ningún título rompa en dos líneas si no fue intencional.
- Confirmar que todas las URLs visibles sean limpias.
- Sustituir el placeholder QR por el destino oficial antes del evento.
""",
    )
    write_markdown(
        "stage-recommendations.md",
        """# Stage Recommendations

## Entrega

- Hablar como fundadora o fundador de una iniciativa, no como instructor.
- Usar pausas deliberadas en WHY y en la pregunta central.
- Mantener el cuerpo erguido y alejarse del atril en las slides 4, 24 y 29.

## Escenografía

- Pantalla limpia, sin overlays innecesarios.
- Clicker confiable.
- Monitor de tiempo visible solo para la persona ponente.

## Narrativa oral

- Evitar explicar botones.
- Repetir la idea central tres veces en distintas formas:
- investigar ya no es solo buscar
- la metodología importa más que la herramienta
- la inteligencia artificial amplifica la capacidad de investigar
""",
    )
    write_markdown(
        "figma-design-spec.md",
        """# Figma Design Spec

## Objetivo

Traducir el keynote a un archivo de Figma Slides o Figma Design manteniendo la misma dirección visual del PowerPoint.

## Dirección visual

- Fondo blanco puro.
- Tipografía sans limpia y dominante.
- Un solo acento rojo oscuro.
- Composición editorial, no UI.

## Componentes base

- Cover title
- Chapter divider
- Narrative statement
- Vertical flow
- Workflow map
- Tool slide template
- Roadmap column
- QR call to action

## Sistema de composición

- Márgenes amplios.
- Alineación izquierda por defecto.
- Slides de capítulos con una sola idea.
- Tool slides con plantilla fija de dos columnas.

## Pendiente antes de publicar

- Crear archivo oficial de Figma Slides en la cuenta del proyecto.
- Reemplazar el placeholder QR con el destino real.
- Validar tipografías disponibles en el entorno Figma del equipo.
""",
    )


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    while prs.slides:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    for item in slides:
        render_slide(prs, item)
    prs.save(OUTPUT_DIR / "research-ai-lab-launch-keynote-v1.pptx")


def main():
    OUTPUT_DIR.mkdir(exist_ok=True)
    PACKAGE_DIR.mkdir(exist_ok=True)
    build_docs()
    build_pptx()


if __name__ == "__main__":
    main()
